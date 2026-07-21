import io
import os
import json
import secrets
import tempfile
from datetime import datetime, timezone

from fastapi import FastAPI, File, UploadFile, HTTPException, Request, BackgroundTasks
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse, FileResponse, JSONResponse
from fastapi.middleware.cors import CORSMiddleware
from fpdf import FPDF
from PIL import Image, ImageDraw, ImageFont
import docx
from docx.table import Table as DocxTable
from docx.text.paragraph import Paragraph as DocxParagraph
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.oxml.ns import qn
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
import qrcode
from qrcode.constants import ERROR_CORRECT_L, ERROR_CORRECT_M, ERROR_CORRECT_Q, ERROR_CORRECT_H

from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded

MAX_FILE_SIZE = 5 * 1024 * 1024  # 5MB
ALLOWED_DOC_EXTENSIONS = {".docx", ".txt", ".pptx"}
API_KEYS_FILE = os.path.join(tempfile.gettempdir(), "api_keys.json")

app = FastAPI()
app.mount("/static", StaticFiles(directory="static"), name="static")

# Allows external developers (using a self-serve API key from /generate-key)
# to call these endpoints from their own website's JavaScript.
app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_methods=["*"],
    allow_headers=["*"],
)


@app.middleware("http")
async def add_security_headers(request: Request, call_next):
    response = await call_next(request)
    # These headers were flagged as missing by securityheaders.com — each
    # closes off a specific class of browser-side attack.
    response.headers["X-Frame-Options"] = "SAMEORIGIN"                    # blocks clickjacking (site being loaded in a hidden iframe)
    response.headers["X-Content-Type-Options"] = "nosniff"                # stops the browser guessing a file's type and running it as something it shouldn't
    response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"  # limits what URL info leaks to other sites when someone clicks a link away from this one
    response.headers["Permissions-Policy"] = "geolocation=(), microphone=(), camera=()"  # explicitly disables browser features this site never uses
    response.headers["Content-Security-Policy"] = (
        "default-src 'self'; "
        "style-src 'self' 'unsafe-inline' https://fonts.googleapis.com; "
        "font-src 'self' https://fonts.gstatic.com; "
        "script-src 'self' 'unsafe-inline'; "
        "img-src 'self' data: blob:; "
        "connect-src 'self';"
    )
    return response


limiter = Limiter(key_func=get_remote_address)
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)


@app.exception_handler(Exception)
async def global_exception_handler(request: Request, exc: Exception):
    return JSONResponse(status_code=500, content={"detail": "Something went wrong on our end. Please try again."})


@app.get("/", response_class=HTMLResponse)
def serve_homepage():
    with open("static/index.html", "r") as f:
        return f.read()


def unique_temp_path(suffix: str) -> str:
    """A fresh, unpredictable filename per request — prevents concurrent
    requests from different users overwriting or reading each other's files."""
    return os.path.join(tempfile.gettempdir(), f"{secrets.token_hex(8)}_{suffix}")


def cleanup_file(path: str):
    try:
        os.remove(path)
    except Exception:
        pass


# ---------------------------------------------------------------------------
# Small shared helpers
# ---------------------------------------------------------------------------

def sanitize_text(text: str) -> str:
    """fpdf2's core fonts only support latin-1, so we replace anything else."""
    return (text or "").encode("latin-1", "replace").decode("latin-1")


def load_watermark_font(size: int):
    font_paths = ["arial.ttf", "Arial.ttf", "C:\\Windows\\Fonts\\arial.ttf", "DejaVuSans-Bold.ttf"]
    for path in font_paths:
        try:
            return ImageFont.truetype(path, size)
        except Exception:
            continue
    return ImageFont.load_default()


def iter_block_items(document):
    """Walk a docx body in document order, yielding Paragraph/Table objects
    so tables show up in the right place instead of being handled separately."""
    for child in document.element.body.iterchildren():
        if child.tag.endswith("}p"):
            yield DocxParagraph(child, document)
        elif child.tag.endswith("}tbl"):
            yield DocxTable(child, document)


def heading_level(style_name: str):
    name = style_name.lower().strip()
    if name.startswith("heading"):
        digits = "".join(ch for ch in name if ch.isdigit())
        return int(digits) if digits else 1
    if name in ("title",):
        return 1
    return 0


def heading_size(level: int) -> int:
    return {1: 18, 2: 16, 3: 14}.get(level, 12)


def docx_align(paragraph) -> str:
    mapping = {
        WD_ALIGN_PARAGRAPH.CENTER: "C",
        WD_ALIGN_PARAGRAPH.RIGHT: "R",
        WD_ALIGN_PARAGRAPH.JUSTIFY: "J",
    }
    return mapping.get(paragraph.alignment, "L")


def pptx_align(paragraph) -> str:
    """PowerPoint uses its own PP_ALIGN enum — distinct from docx's
    WD_ALIGN_PARAGRAPH, even though the values look similar."""
    mapping = {
        PP_ALIGN.CENTER: "C",
        PP_ALIGN.RIGHT: "R",
        PP_ALIGN.JUSTIFY: "J",
    }
    try:
        return mapping.get(paragraph.alignment, "L")
    except Exception:
        return "L"


def run_color(run):
    """Returns (r,g,b) if the run has an explicit color set, else None."""
    try:
        color = run.font.color
        if color and color.type is not None and color.rgb is not None:
            rgb = color.rgb
            return (rgb[0], rgb[1], rgb[2])
    except Exception:
        pass
    return None


def emu_to_mm(emu) -> float:
    return float(emu) / 36000.0


def length_to_mm(length) -> float:
    """Convert a docx Length (EMU-based) to millimeters."""
    if length is None:
        return 0.0
    return emu_to_mm(length)


def extract_docx_paragraph_images(para: DocxParagraph, document):
    """Finds any inline images referenced inside a paragraph's runs and
    returns a list of raw image bytes, in the order they appear."""
    images = []
    for run in para.runs:
        blips = run._element.findall(".//" + qn("a:blip"))
        for blip in blips:
            r_id = blip.get(qn("r:embed"))
            if not r_id:
                continue
            try:
                image_part = document.part.related_parts[r_id]
                images.append(image_part.blob)
            except Exception:
                continue
    return images


def draw_pdf_image(pdf: FPDF, image_bytes: bytes, max_width_mm: float):
    """Places an image at the current cursor position, scaled to fit
    max_width_mm while keeping its aspect ratio, then advances the cursor."""
    try:
        pil_img = Image.open(io.BytesIO(image_bytes))
        px_w, px_h = pil_img.size
    except Exception:
        return
    width_mm = min(max_width_mm, (px_w / 96.0) * 25.4)
    height_mm = width_mm * (px_h / px_w) if px_w else 0
    if pdf.get_y() + height_mm > pdf.h - pdf.b_margin:
        pdf.add_page()
    x = pdf.get_x()
    y = pdf.get_y()
    try:
        pdf.image(io.BytesIO(image_bytes), x=x, y=y, w=width_mm, h=height_mm)
        pdf.set_y(y + height_mm + 4)
    except Exception:
        pass


def cell_shading_color(cell):
    """Reads a table cell's background fill color (w:shd) if one is set."""
    try:
        tc_pr = cell._tc.find(qn("w:tcPr"))
        if tc_pr is None:
            return None
        shd = tc_pr.find(qn("w:shd"))
        if shd is None:
            return None
        fill = shd.get(qn("w:fill"))
        if not fill or fill.lower() == "auto":
            return None
        fill = fill.lstrip("#")
        return tuple(int(fill[i:i + 2], 16) for i in (0, 2, 4))
    except Exception:
        return None


# ---------------------------------------------------------------------------
# Compress image (unchanged)
# ---------------------------------------------------------------------------

@app.post("/compress-image")
@limiter.limit("10/minute")
async def compress_image(request: Request, background_tasks: BackgroundTasks, file: UploadFile = File(...), quality: int = 50):
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="File must be an image.")
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 5MB.")
    try:
        image = Image.open(io.BytesIO(contents))
        image = image.convert("RGB")
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted image file.")
    safe_quality = max(1, min(95, quality))
    temp_path = unique_temp_path("compressed_output.jpg")
    image.save(temp_path, format="JPEG", quality=safe_quality, optimize=True)
    background_tasks.add_task(cleanup_file, temp_path)
    return FileResponse(temp_path, media_type="image/jpeg", filename="compressed.jpg", background=background_tasks)


# ---------------------------------------------------------------------------
# Document -> PDF, with real formatting (bold/italic/underline, headings,
# bullet lists, and basic tables for .docx; titles + indented bullets per
# slide for .pptx; plain paragraphs for .txt)
# ---------------------------------------------------------------------------

def render_docx_paragraph(pdf: FPDF, para: DocxParagraph, document):
    style_name = para.style.name if para.style else ""
    level = heading_level(style_name)
    is_list = "list" in style_name.lower() or "bullet" in style_name.lower()
    text = para.text.strip()

    # Inline images live inside runs even when there's no text alongside them.
    images = extract_docx_paragraph_images(para, document)

    if not text and not images:
        return

    align = docx_align(para)
    base_indent = length_to_mm(para.paragraph_format.left_indent)
    indent = max(base_indent, 6 if is_list else 0)
    pdf.set_x(pdf.l_margin + indent)

    space_before = length_to_mm(para.paragraph_format.space_before)
    if space_before:
        pdf.ln(min(space_before, 8))

    if text:
        if is_list:
            pdf.set_font("Helvetica", "", 12)
            pdf.write(7, "-  ")

        runs = para.runs if para.runs else None
        line_height = 9 if level else 7

        if not runs:
            size = heading_size(level) if level else 12
            style = "B" if level else ""
            pdf.set_font("Helvetica", style, size)
            full_text = sanitize_text(text)
            if align in ("C", "R", "J"):
                usable_width = pdf.w - pdf.l_margin - pdf.r_margin - indent
                pdf.multi_cell(usable_width, line_height, full_text, align=align)
            else:
                pdf.write(line_height, full_text)
                pdf.ln(line_height + (1 if level else 0))
        else:
            # Gather each run's own style first, so alignment and formatting
            # can both be honored instead of one overriding the other.
            run_specs = []
            for run in runs:
                run_text = sanitize_text(run.text)
                if not run_text:
                    continue
                style = ""
                if run.bold or level:
                    style += "B"
                if run.italic:
                    style += "I"
                if run.underline:
                    style += "U"
                size = heading_size(level) if level else 12
                color = run_color(run)
                run_specs.append((run_text, style, size, color))

            if align in ("C", "R") and run_specs:
                # write() always flows left-to-right, so for center/right we
                # measure the full line first and offset the cursor before
                # drawing each run — that way per-run bold/italic/color still
                # comes through instead of being dropped for alignment's sake.
                total_width = 0
                for run_text, style, size, _ in run_specs:
                    pdf.set_font("Helvetica", style, size)
                    total_width += pdf.get_string_width(run_text)
                usable_width = pdf.w - pdf.l_margin - pdf.r_margin - indent
                offset = (usable_width - total_width) if align == "R" else (usable_width - total_width) / 2
                pdf.set_x(pdf.l_margin + indent + max(offset, 0))
            # Note: true justify (even word-spacing) across mixed-style runs
            # isn't something fpdf2 supports directly, so justified paragraphs
            # with multiple runs fall back to left flow here rather than
            # losing bold/italic/color — a deliberate, documented trade-off.

            for run_text, style, size, color in run_specs:
                pdf.set_font("Helvetica", style, size)
                if color:
                    pdf.set_text_color(*color)
                pdf.write(line_height, run_text)
                if color:
                    pdf.set_text_color(0, 0, 0)
            pdf.ln(line_height + (1 if level else 0))

    for image_bytes in images:
        usable_width = pdf.w - pdf.l_margin - pdf.r_margin - indent
        pdf.set_x(pdf.l_margin + indent)
        draw_pdf_image(pdf, image_bytes, usable_width)


def render_docx_table(pdf: FPDF, table: DocxTable):
    col_count = max(len(table.columns), 1)
    usable_width = pdf.w - pdf.l_margin - pdf.r_margin
    col_width = usable_width / col_count
    row_height = 7

    pdf.ln(2)
    for row_index, row in enumerate(table.rows):
        x_start = pdf.get_x()
        y_start = pdf.get_y()
        max_y = y_start
        is_header = row_index == 0
        for i, cell in enumerate(row.cells):
            cell_text = sanitize_text(cell.text.strip())
            pdf.set_xy(x_start + i * col_width, y_start)
            pdf.set_font("Helvetica", "B" if is_header else "", 10)

            fill = cell_shading_color(cell)
            if fill:
                pdf.set_fill_color(*fill)
            elif is_header:
                pdf.set_fill_color(230, 230, 230)
            else:
                pdf.set_fill_color(255, 255, 255)

            pdf.multi_cell(col_width, row_height, cell_text, border=1, fill=True)
            if pdf.get_y() > max_y:
                max_y = pdf.get_y()
        pdf.set_xy(x_start, max_y)
    pdf.ln(4)


def build_pdf_from_docx(contents: bytes) -> FPDF:
    document = docx.Document(io.BytesIO(contents))
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()

    found_any = False
    for block in iter_block_items(document):
        if isinstance(block, DocxParagraph):
            if block.text.strip() or extract_docx_paragraph_images(block, document):
                found_any = True
            render_docx_paragraph(pdf, block, document)
        elif isinstance(block, DocxTable):
            if block.rows:
                found_any = True
            render_docx_table(pdf, block)

    if not found_any:
        raise ValueError("empty document")
    return pdf


def build_pdf_from_pptx(contents: bytes) -> FPDF:
    prs = Presentation(io.BytesIO(contents))
    slide_w_mm = emu_to_mm(prs.slide_width)
    slide_h_mm = emu_to_mm(prs.slide_height)

    # NOTE: format=(w,h) is used exactly as given only when orientation="P" —
    # passing "L" here would make fpdf2 swap the two values, which would
    # corrupt every shape's x/y position on the far more common landscape slides.
    pdf = FPDF(orientation="P", unit="mm", format=(slide_w_mm, slide_h_mm))
    pdf.set_auto_page_break(auto=False)

    found_any = False
    for slide in prs.slides:
        pdf.add_page()

        # Sort so pictures are drawn first and text sits on top of them.
        shapes = sorted(slide.shapes, key=lambda s: 0 if s.shape_type == MSO_SHAPE_TYPE.PICTURE else 1)

        for shape in shapes:
            x_mm = emu_to_mm(shape.left) if shape.left is not None else 0
            y_mm = emu_to_mm(shape.top) if shape.top is not None else 0
            w_mm = emu_to_mm(shape.width) if shape.width is not None else slide_w_mm
            h_mm = emu_to_mm(shape.height) if shape.height is not None else 10

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    pdf.image(io.BytesIO(shape.image.blob), x=x_mm, y=y_mm, w=w_mm, h=h_mm)
                    found_any = True
                except Exception:
                    pass
                continue

            if not shape.has_text_frame:
                continue

            is_title = shape == slide.shapes.title
            cursor_y = y_mm
            for para in shape.text_frame.paragraphs:
                text = "".join(run.text for run in para.runs).strip()
                if not text:
                    continue
                found_any = True

                indent = 0 if is_title else min(para.level, 4) * 5
                align = pptx_align(para)
                bullet = "" if is_title else "-  "

                pdf.set_xy(x_mm + indent, cursor_y)
                first_run = para.runs[0] if para.runs else None
                size = 12
                if first_run is not None and first_run.font.size is not None:
                    size = first_run.font.size.pt
                elif is_title:
                    size = 24
                style = "B" if (is_title or (first_run and first_run.font.bold)) else ""
                if first_run and first_run.font.italic:
                    style += "I"

                pdf.set_font("Helvetica", style, size)
                color = run_color(first_run) if first_run else None
                if color:
                    pdf.set_text_color(*color)

                line_height = max(size * 0.45, 5)
                pdf.multi_cell(max(w_mm - indent, 10), line_height, sanitize_text(bullet + text), align=align)
                if color:
                    pdf.set_text_color(0, 0, 0)

                cursor_y = pdf.get_y() + 1
                if cursor_y > y_mm + h_mm + 20:
                    break

    if not found_any:
        raise ValueError("empty document")
    return pdf


def build_pdf_from_txt(contents: bytes) -> FPDF:
    text_content = contents.decode("utf-8", errors="ignore")
    lines = [line.strip() for line in text_content.splitlines() if line.strip()]
    if not lines:
        raise ValueError("empty document")

    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=15)
    pdf.add_page()
    pdf.set_font("Helvetica", "", 12)
    for line in lines:
        pdf.multi_cell(0, 8, sanitize_text(line))
        pdf.ln(1)
    return pdf


@app.post("/convert-document-to-pdf")
@limiter.limit("10/minute")
async def convert_document_to_pdf(request: Request, background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    filename = file.filename or ""
    ext = os.path.splitext(filename)[1].lower()

    if ext not in ALLOWED_DOC_EXTENSIONS:
        raise HTTPException(status_code=400, detail="Unsupported file type. Please upload a .docx, .txt, or .pptx file.")

    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 5MB.")

    try:
        if ext == ".docx":
            pdf = build_pdf_from_docx(contents)
        elif ext == ".pptx":
            pdf = build_pdf_from_pptx(contents)
        else:
            pdf = build_pdf_from_txt(contents)
    except ValueError:
        raise HTTPException(status_code=400, detail="The document appears to be empty.")
    except Exception:
        raise HTTPException(status_code=400, detail="Could not read this document. It may be corrupted or invalid.")

    temp_path = unique_temp_path("document_converted.pdf")
    pdf.output(temp_path)
    background_tasks.add_task(cleanup_file, temp_path)
    return FileResponse(temp_path, media_type="application/pdf", filename="converted_document.pdf", background=background_tasks)


# ---------------------------------------------------------------------------
# Image -> PDF (unchanged)
# ---------------------------------------------------------------------------

@app.post("/convert-image-to-pdf")
@limiter.limit("10/minute")
async def convert_image_to_pdf(request: Request, background_tasks: BackgroundTasks, file: UploadFile = File(...)):
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="File must be an image.")
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 5MB.")
    try:
        image = Image.open(io.BytesIO(contents))
        image = image.convert("RGB")
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted image file.")
    temp_path = unique_temp_path("converted.pdf")
    image.save(temp_path)
    background_tasks.add_task(cleanup_file, temp_path)
    return FileResponse(temp_path, media_type="application/pdf", filename="converted.pdf", background=background_tasks)


# ---------------------------------------------------------------------------
# Watermark image — single placement OR tiled/repeated pattern
# ---------------------------------------------------------------------------

@app.post("/watermark-image")
@limiter.limit("10/minute")
async def watermark_image(
    request: Request,
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    text: str = "SAMPLE",
    mode: str = "single",          # "single" or "tiled"
    position: str = "bottom-right",  # used only when mode == "single"
    opacity: int = 180,
    angle: int = -30,              # used only when mode == "tiled"
):
    if not file.content_type or not file.content_type.startswith("image/"):
        raise HTTPException(status_code=400, detail="File must be an image.")
    contents = await file.read()
    if len(contents) > MAX_FILE_SIZE:
        raise HTTPException(status_code=413, detail="File too large. Max size is 5MB.")
    try:
        image = Image.open(io.BytesIO(contents)).convert("RGBA")
    except Exception:
        raise HTTPException(status_code=400, detail="Invalid or corrupted image file.")

    safe_opacity = max(0, min(255, opacity))
    overlay = Image.new("RGBA", image.size, (255, 255, 255, 0))
    draw = ImageDraw.Draw(overlay)

    if mode == "tiled":
        tile_font_size = max(18, image.width // 20)
        tile_font = load_watermark_font(tile_font_size)

        bbox = draw.textbbox((0, 0), text, font=tile_font)
        tw, th = bbox[2] - bbox[0], bbox[3] - bbox[1]
        pad = 40
        tile = Image.new("RGBA", (tw + pad * 2, th + pad * 2), (255, 255, 255, 0))
        tile_draw = ImageDraw.Draw(tile)
        tile_draw.text((pad, pad), text, font=tile_font, fill=(255, 255, 255, safe_opacity))
        rotated_tile = tile.rotate(angle, expand=True)
        tile_w, tile_h = rotated_tile.size

        spacing_x = tile_w + 30
        spacing_y = tile_h + 30
        for y in range(-tile_h, image.height + tile_h, spacing_y):
            for x in range(-tile_w, image.width + tile_w, spacing_x):
                overlay.paste(rotated_tile, (x, y), rotated_tile)
    else:
        font_size = max(24, image.width // 15)
        font = load_watermark_font(font_size)
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        margin = 20

        positions = {
            "top-left": (margin, margin),
            "top-right": (image.width - text_width - margin, margin),
            "bottom-left": (margin, image.height - text_height - margin),
            "bottom-right": (image.width - text_width - margin, image.height - text_height - margin),
            "center": ((image.width - text_width) // 2, (image.height - text_height) // 2),
        }
        xy = positions.get(position, positions["bottom-right"])
        draw.text(xy, text, font=font, fill=(255, 255, 255, safe_opacity), stroke_width=3, stroke_fill=(0, 0, 0, safe_opacity))

    watermarked = Image.alpha_composite(image, overlay).convert("RGB")
    temp_path = unique_temp_path("watermarked.jpg")
    watermarked.save(temp_path, format="JPEG")
    background_tasks.add_task(cleanup_file, temp_path)
    return FileResponse(temp_path, media_type="image/jpeg", filename="watermarked.jpg", background=background_tasks)


# ---------------------------------------------------------------------------
# QR code generator
# ---------------------------------------------------------------------------

QR_ERROR_LEVELS = {
    "L": ERROR_CORRECT_L,  # ~7% of the code can be damaged/covered and still scan
    "M": ERROR_CORRECT_M,  # ~15%
    "Q": ERROR_CORRECT_Q,  # ~25%
    "H": ERROR_CORRECT_H,  # ~30% — most reliable, but denser-looking pattern
}
MAX_QR_TEXT_LENGTH = 1000


@app.post("/generate-qr-code")
@limiter.limit("10/minute")
async def generate_qr_code(request: Request, background_tasks: BackgroundTasks, text: str = "", size: int = 8, error_correction: str = "M"):
    clean_text = (text or "").strip()
    if not clean_text:
        raise HTTPException(status_code=400, detail="Please enter some text or a link to encode.")
    if len(clean_text) > MAX_QR_TEXT_LENGTH:
        raise HTTPException(status_code=400, detail=f"Text is too long (max {MAX_QR_TEXT_LENGTH} characters).")

    safe_size = max(2, min(20, size))
    level = QR_ERROR_LEVELS.get(error_correction.upper(), ERROR_CORRECT_M)

    qr = qrcode.QRCode(
        version=None,
        error_correction=level,
        box_size=safe_size,
        border=4,
    )
    qr.add_data(clean_text)
    qr.make(fit=True)
    img = qr.make_image(fill_color="black", back_color="white").convert("RGB")

    temp_path = unique_temp_path("qrcode.png")
    img.save(temp_path, format="PNG")
    background_tasks.add_task(cleanup_file, temp_path)
    return FileResponse(temp_path, media_type="image/png", filename="qrcode.png", background=background_tasks)


# ---------------------------------------------------------------------------
# Self-serve API key generation (informational only — not enforced on
# endpoints). Keys are logged to a JSON file so there's a record of who
# generated one, even though nothing currently checks for it.
# ---------------------------------------------------------------------------

def load_issued_keys():
    if not os.path.exists(API_KEYS_FILE):
        return []
    try:
        with open(API_KEYS_FILE, "r") as f:
            return json.load(f)
    except Exception:
        return []


def save_issued_keys(keys):
    with open(API_KEYS_FILE, "w") as f:
        json.dump(keys, f, indent=2)


@app.post("/generate-key")
@limiter.limit("5/minute")
async def generate_key(request: Request, email: str = ""):
    clean_email = (email or "").strip()[:100]
    if "@" not in clean_email or "." not in clean_email.split("@")[-1] or len(clean_email) < 5:
        raise HTTPException(status_code=400, detail="Please enter a valid email address.")

    new_key = f"fpapi_{secrets.token_hex(16)}"

    keys = load_issued_keys()
    keys.append({
        "email": clean_email,
        "key": new_key,
        "date": datetime.now(timezone.utc).strftime("%Y-%m-%d %H:%M UTC"),
    })
    save_issued_keys(keys)

    return {"email": clean_email, "api_key": new_key}